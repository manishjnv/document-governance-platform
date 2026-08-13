"""Report/export/compare tests (MITRE Phase 4). Assessments are seeded
directly with stored summary/technique_results — no pipeline, no LLM."""

import io
import uuid
from datetime import datetime, timedelta, timezone

import pytest
from httpx import ASGITransport, AsyncClient
from openpyxl import load_workbook

from app.auth import create_access_token
from app.mitre.report import _guard, build_html_report, build_xlsx_export
from app.mitre.service import compare_assessments
from app.models.mitre_assessment import MitreAssessment
from app.models.mitre_use_case import MitreUseCase
from app.models.organization import Organization
from app.models.user import User
from main import app

XSS = "<script>alert(1)</script>"
FORMULA = "=2+2"


@pytest.fixture
async def client():
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://testserver") as c:
        yield c


async def _make_user(db_session, *, role="admin"):
    org = Organization(org_id=uuid.uuid4(), name=f"org-{uuid.uuid4()}")
    user = User(user_id=uuid.uuid4(), org_id=org.org_id, email=f"{uuid.uuid4()}@x.com")
    db_session.add_all([org, user])
    await db_session.commit()
    token, _ = create_access_token(
        user_id=user.user_id, email=user.email, org_id=org.org_id, role=role
    )
    return org, user, {"Authorization": f"Bearer {token}"}


def _gap(tid, name, rank=1, state="not_covered", threat_relevance=("Banking",),
          crown_jewel_relevant=True):
    return {
        "technique_id": tid, "name": name, "domain": "enterprise", "state": state,
        "tier": 2, "tactics": ["TA0002"], "feasibility": "short", "via": "Sysmon",
        "category": "registry", "hint": "build the detection now", "rank": rank,
        "threat_relevance": list(threat_relevance) if threat_relevance else None,
        "crown_jewel_relevant": crown_jewel_relevant,
    }


def _tactic(strict_pct):
    return {
        "id": "TA0002", "shortname": "execution", "name": "Execution",
        "covered": 1, "partial": 1, "not_covered": 1, "not_applicable": 1,
        "applicable": 3, "strict_pct": strict_pct, "weighted_pct": strict_pct + 10,
    }


def _summary(strict_pct=33.3):
    gap = _gap("T1112", "Modify Registry")
    return {
        "overall": {"covered": 1, "partial": 1, "not_covered": 1, "not_applicable": 1,
                    "applicable": 3, "strict_pct": strict_pct, "weighted_pct": 50.0},
        "domains": {"enterprise": {"covered": 1, "partial": 1, "not_covered": 1,
                                   "not_applicable": 1, "applicable": 3,
                                   "strict_pct": strict_pct, "weighted_pct": 50.0,
                                   "tactics": [_tactic(strict_pct)]}},
        "assumptions": [f"seeded assumption {XSS}"],
        "gaps": [gap],
        "roadmap": {"short": [gap], "mid": [], "long": []},
        "narrative": {"executive_summary": "Coverage is low.",
                      "gap_recommendations": {"T1112": "Build a Sysmon registry detection."},
                      "roadmap_prose": {"short": "s", "mid": "m", "long": "l"},
                      "generated_by": "template", "model_used": None},
        "not_applicable": [{"technique_id": "T1200", "domain": "enterprise",
                            "reason": "accepted risk, physical controls"}],
        "counts": {"use_cases": 2, "customer_tagged": 2, "ai_tagged": 0,
                   "unmapped": 0, "invalid": 0},
        "applicable_domains": ["enterprise"],
    }


def _results(states: dict):
    return [
        {"technique_id": tid, "domain": "enterprise", "tactics": ["TA0002"],
         "state": state, "na_reason": "reason" if state == "not_applicable" else None,
         "use_case_refs": []}
        for tid, state in states.items()
    ]


async def _seed(db_session, org, user, *, status="completed", strict_pct=33.3,
                 results=None, uc_names=(), name="Seeded assessment",
                 customer=None, completed_at=None):
    params = {"thresholds": {"confidence_covered": 0.7}, "models_used": {}}
    if customer is not None:
        params["customer"] = customer
    assessment = MitreAssessment(
        assessment_id=uuid.uuid4(),
        org_id=org.org_id,
        name=name,
        status=status,
        attack_version="19.1",
        summary=_summary(strict_pct) if status == "completed" else None,
        technique_results=results
        if results is not None
        else _results({"T1059.001": "covered", "T1003.001": "partial",
                       "T1112": "not_covered", "T1200": "not_applicable"}),
        completed_at=(
            (completed_at or datetime.now(timezone.utc)) if status == "completed" else None
        ),
        error_message="boom" if status == "failed" else None,
        created_by=user.user_id,
        params=params,
    )
    db_session.add(assessment)
    for i, name in enumerate(uc_names, start=1):
        db_session.add(
            MitreUseCase(
                assessment_id=assessment.assessment_id,
                org_id=org.org_id,
                row_ref=f"s:{i}",
                name=name,
                description=f"desc {name}",
                enabled=True,
                mappings=[{"technique_id": "T1059.001", "source": "customer", "confidence": 1.0}],
                mapping_status="customer_tagged",
            )
        )
    await db_session.commit()
    return assessment


# ---------------------------------------------------------------------------
# pure-function tests
# ---------------------------------------------------------------------------


def test_formula_guard():
    assert _guard("=2+2") == "'=2+2"
    assert _guard("+SUM(A1)") == "'+SUM(A1)"
    assert _guard("-1") == "'-1"
    assert _guard("@cmd") == "'@cmd"
    assert _guard("normal") == "normal"
    assert _guard(5) == 5
    assert _guard(None) is None


# --- Phase A10 piece 3: coverage by log source ---

def test_log_source_coverage_groups_by_normalized_log_source():
    from app.mitre import attack_data
    from app.mitre.report_common import OTHER_LOG_SOURCE, compute_log_source_coverage

    use_cases = [
        {"row_ref": "s:1", "log_source": "Sysmon",
         "mappings": [{"technique_id": "T1059.001"}]},
        {"row_ref": "s:2", "log_source": "sysmon",  # case variant -> same group
         "mappings": [{"technique_id": "T1112"}]},
        {"row_ref": "s:3", "log_source": "CloudTrail",
         "mappings": [{"technique_id": "T1078"}]},
        {"row_ref": "s:4", "log_source": None, "mappings": []},
    ]
    technique_results = _results(
        {"T1059.001": "covered", "T1112": "partial", "T1078": "not_covered"}
    )
    groups = compute_log_source_coverage(use_cases, technique_results, attack_data.DEFAULT)
    by_name = {g["log_source"]: g for g in groups}

    assert by_name["Sysmon"]["rule_count"] == 2
    assert by_name["Sysmon"]["techniques_covered"] == 2
    assert {t["technique_id"] for t in by_name["Sysmon"]["techniques"]} == {
        "T1059.001", "T1112",
    }
    assert "Execution" in by_name["Sysmon"]["tactics"]
    assert by_name["Sysmon"]["row_refs"] == ["s:1", "s:2"]
    assert by_name["CloudTrail"]["rule_count"] == 1

    # rules with no log_source never get dropped
    assert OTHER_LOG_SOURCE in by_name
    assert by_name[OTHER_LOG_SOURCE]["rule_count"] == 1
    # "Other" sorts last regardless of count
    assert groups[-1]["log_source"] == OTHER_LOG_SOURCE


def test_log_source_coverage_ignores_unresolvable_technique_ids():
    from app.mitre import attack_data
    from app.mitre.report_common import compute_log_source_coverage

    use_cases = [
        {"row_ref": "s:1", "log_source": "Sysmon",
         "mappings": [{"technique_id": "T9999.999"}]},
    ]
    groups = compute_log_source_coverage(use_cases, [], attack_data.DEFAULT)
    assert groups[0]["techniques"] == []
    assert groups[0]["techniques_covered"] == 0


@pytest.mark.asyncio
async def test_html_report_escapes_untrusted_strings(db_session):
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user, uc_names=[f"Rule {XSS}"])
    use_cases = [{"row_ref": "s:1", "name": f"Rule {XSS}", "description": None,
                  "log_source": None, "enabled": True, "mappings": [],
                  "mapping_status": "customer_tagged"}]
    html = build_html_report(assessment, use_cases)
    assert XSS not in html                       # raw payload never present
    assert "&lt;script&gt;" in html              # escaped twice: rule + assumption
    assert "33.3%" in html                       # numbers from stored summary
    assert "Build a Sysmon registry detection." in html
    assert "accepted risk, physical controls" in html
    # Phase 14e structure: TOC with real page numbers + per-gap why line
    assert "target-counter(attr(href), page)" in html
    assert "Why it's a gap:" in html
    assert "Top 5 things to fix first" in html


@pytest.mark.asyncio
async def test_html_report_roadmap_index_and_register_dedup(db_session):
    """Phase A9: the roadmap section is a compact per-bucket index (ID, Name,
    Priority, cross-ref) that points into the gap register; the gap's full
    narrative (why/recommendation) prints exactly once, in the register —
    not repeated inline in the roadmap."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    html = build_html_report(assessment, [])

    assert '<h2 id="roadmap"' in html
    assert '<h2 id="gapreg"' in html
    assert "Gap register" in html
    roadmap_start = html.index('<h2 id="roadmap"')
    register_start = html.index('<h2 id="gapreg"')
    assert roadmap_start < register_start
    roadmap_slice = html[roadmap_start:register_start]
    register_slice = html[register_start:]

    # index table: technique id + cross-ref present, full narrative absent
    assert "T1112" in roadmap_slice
    assert "<a class='xref' href='#g-T1112'></a>" in roadmap_slice
    assert "Build a Sysmon registry detection." not in roadmap_slice
    assert "Why it's a gap:" not in roadmap_slice

    # single home: the full narrative and its anchor appear exactly once
    assert "Build a Sysmon registry detection." in register_slice
    assert register_slice.count("id='g-T1112'") == 1


@pytest.mark.asyncio
async def test_html_report_gaps_scope_keeps_roadmap_and_register(db_session):
    """Phase A9: the gaps tab scope keeps both the roadmap index and the
    (untouched) gap register; the coverage tab scope keeps neither."""
    org, user, _headers = await _make_user(db_session, role="viewer")
    assessment = await _seed(db_session, org, user)

    gaps_html = build_html_report(assessment, [], scope="gaps")
    assert '<h2 id="roadmap"' in gaps_html
    assert "Gap register" in gaps_html

    coverage_html = build_html_report(assessment, [], scope="coverage")
    assert '<h2 id="roadmap"' not in coverage_html
    assert "Gap register" not in coverage_html
    assert "Coverage by attack stage" in coverage_html


@pytest.mark.asyncio
async def test_html_report_page_break_audit(db_session):
    """Phase A11 piece 3: hard page breaks survive ONLY at genuine PART
    boundaries (cover->executive->detailed->appendix); sub-sections inside
    a part (roadmap, inside "detailed") flow continuously so content isn't
    stranded on a near-empty page."""
    org, user, _headers = await _make_user(db_session, role="viewer")
    assessment = await _seed(db_session, org, user)
    html = build_html_report(assessment, [])

    def has_break(anchor_id):
        return f'id="{anchor_id}" class="page-break"' in html

    assert has_break("exec")       # cover -> executive
    assert has_break("tactics")    # executive -> detailed
    assert has_break("register")   # detailed -> appendix
    assert not has_break("roadmap")  # sub-section within "detailed" -- flows


@pytest.mark.asyncio
async def test_xlsx_formula_injection_guard(db_session):
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    use_cases = [{"row_ref": "s:1", "name": FORMULA, "description": "=HYPERLINK(evil)",
                  "logic": "=cmd|'/C calc'!A0", "log_source": None, "enabled": True,
                  "mappings": [], "mapping_status": "customer_tagged"}]
    wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, use_cases)))
    assert set(wb.sheetnames) == {
        "Read Me", "Summary", "Coverage by Tactic", "Technique Tracker",
        "Use-Case Mappings", "Coverage by Log Source", "Not Applicable", "Assumptions",
    }
    assert wb.sheetnames[0] == "Read Me"                # Phase 14c guide sheet
    ws = wb["Use-Case Mappings"]
    assert ws["B2"].value == "'" + FORMULA
    assert ws["I2"].value == "'=HYPERLINK(evil)"
    assert ws["J1"].value == "Logic"                    # Phase 7 column
    assert ws["J2"].value == "'=cmd|'/C calc'!A0"       # logic cell guarded


@pytest.mark.asyncio
async def test_xlsx_tracker_structure(db_session):
    """Phase A9: the merged Technique Tracker sheet — exact header order, no
    interleaved section-header rows, one row per APPLICABLE technique (N/A
    excluded), covered rows leave gap-only columns blank, a gap row is fully
    populated including the new Threat match/Crown jewel/Roadmap bucket/
    Owner-Status-Target-date-Notes columns."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    use_cases = [
        {"row_ref": "s:10", "name": "Rule ten", "description": None, "logic": None,
         "log_source": None, "enabled": True, "mappings": [],
         "mapping_status": "unmapped"},
        {"row_ref": "s:2", "name": "PS rule", "description": None, "logic": None,
         "log_source": None, "enabled": True,
         "mappings": [{"technique_id": "T1059.001", "source": "customer",
                       "confidence": 1.0}],
         "mapping_status": "customer_tagged"},
    ]
    wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, use_cases)))

    readme = wb["Read Me"]
    texts = [str(c.value) for row in readme.iter_rows() for c in row if c.value]
    assert any("Is 33.3% bad?" in t for t in texts)
    assert any("What each sheet contains" in t for t in texts)
    assert any("Technique Tracker" in t for t in texts)

    tracker = wb["Technique Tracker"]
    headers = [c.value for c in tracker[1]]
    assert headers == [
        "Technique ID", "Name", "Tactic(s)", "Domain", "State", "Why", "Strength",
        "Priority", "Threat match", "Crown jewel", "Feasibility", "Roadmap bucket",
        "Recommendation", "Log fields needed",
        "Reference KQL (illustrative — tune before use)", "Via", "Owner", "Status",
        "Target date", "Notes",
    ]
    # no interleaved section-header rows: every data row's first cell is a
    # real technique id from the applicable set (N/A technique T1200 excluded)
    ids = [r[0].value for r in tracker.iter_rows(min_row=2)]
    assert all(str(i).startswith("T") for i in ids)
    assert len(ids) == 3 and "T1200" not in ids
    assert set(ids) == {"T1059.001", "T1003.001", "T1112"}

    rows = {r[0].value: r for r in tracker.iter_rows(min_row=2)}
    covered = rows["T1059.001"]
    assert covered[1].value == "PowerShell"             # name from the pinned index
    assert covered[2].value == "Execution"              # tactic name, not TA0002
    assert "Covered by your rule 'PS rule'" in covered[5].value
    # covered row: gap-only columns blank (no gap entry for a covered technique)
    for col in (7, 8, 9, 10, 11, 12, 13, 14):
        assert covered[col].value in (None, "")

    gap_row = rows["T1112"]
    assert gap_row[4].value == "No rule detects this"   # State (plain words)
    assert "maps to this technique" in gap_row[5].value  # Why
    assert gap_row[7].value == 2                        # Priority: numeric P2
    assert gap_row[7].number_format == '"P"0'
    assert gap_row[8].value == "Banking"                # Threat match
    assert gap_row[9].value == "Yes"                    # Crown jewel
    assert gap_row[10].value == "Short term (0–3 mo)"   # Feasibility
    assert gap_row[11].value == "Short"                 # Roadmap bucket VALUE
    assert gap_row[12].value == "Build a Sysmon registry detection."  # Recommendation
    assert "your query needs:" in gap_row[13].value      # Log fields needed
    assert "Windows Registry Key Modification" in gap_row[13].value
    # Reference KQL (2026-08-14): single-word via -> skeleton with the
    # never-fires + false-positive discipline header
    assert "REFERENCE ONLY" in gap_row[14].value
    assert "Sysmon | take 10" in gap_row[14].value
    assert gap_row[15].value == "Sysmon"                # Via
    # blank customer-tracking columns
    for col in (16, 17, 18, 19):
        assert gap_row[col].value in (None, "")

    ucs = wb["Use-Case Mappings"]
    assert ucs["A2"].value == "s:2"                     # numeric sort: 2 before 10
    assert ucs["A3"].value == "s:10"
    assert ucs["D2"].value == "You tagged this"         # plain-words status
    assert ucs["D3"].value == "Could not be mapped"

    summary_metrics = [r[0].value for r in wb["Summary"].iter_rows(min_row=2)]
    assert "Coverage %" in summary_metrics
    assert "Strict coverage %" not in summary_metrics


def test_xlsx_tracker_formula_guard():
    """Phase A9: the Tracker's merged Recommendation column carries the same
    attacker-controlled narrative text the old Gaps & Recommendations sheet
    did — the formula-injection guard must still apply."""
    class A:
        pass

    a = A()
    a.assessment_id = uuid.uuid4()
    a.name = "x"
    a.completed_at = datetime.now(timezone.utc)
    a.attack_version = "19.1"
    a.summary = _summary()
    a.summary["narrative"]["gap_recommendations"]["T1112"] = FORMULA
    a.technique_results = _results({"T1112": "not_covered"})
    a.params = {"thresholds": {"confidence_covered": 0.7}, "models_used": {}}

    wb = load_workbook(io.BytesIO(build_xlsx_export(a, [])))
    tracker = wb["Technique Tracker"]
    row = next(r for r in tracker.iter_rows(min_row=2) if r[0].value == "T1112")
    assert row[12].value == "'" + FORMULA               # Recommendation, guarded


@pytest.mark.asyncio
async def test_xlsx_a11_uniform_header_fill(db_session):
    """Phase A11 piece 1: every sheet's header row shares ONE consistent
    branded fill + white bold font (previously bold-but-unfilled on most
    sheets; Summary's mini-table headers were bold-only too, only its
    section-title bars were brand-filled)."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, [])))

    brand_rgb, white_rgb = "00341954", "00FFFFFF"  # 2026-08-14 purple restyle
    for name in (
        "Read Me", "Coverage by Tactic", "Technique Tracker", "Use-Case Mappings",
        "Coverage by Log Source", "Not Applicable", "Assumptions",
    ):
        header_cell = wb[name].cell(row=1, column=1)
        assert header_cell.fill.fgColor.rgb == brand_rgb, name
        assert header_cell.font.color.rgb == white_rgb, name

    # Summary's mini-table headers get the same treatment (previously bold-only)
    ws_sum = wb["Summary"]
    metric_header = next(r for r in ws_sum.iter_rows() if r[0].value == "Metric")
    assert metric_header[0].fill.fgColor.rgb == brand_rgb
    assert metric_header[0].font.color.rgb == white_rgb


@pytest.mark.asyncio
async def test_xlsx_scope_pruning(db_session):
    """Phase A9: coverage and gaps per-tab downloads both keep the merged
    Tracker sheet (it now carries both roles); assumptions is unaffected."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)

    coverage_wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, [], scope="coverage")))
    assert set(coverage_wb.sheetnames) == {"Coverage by Tactic", "Technique Tracker"}

    gaps_wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, [], scope="gaps")))
    assert set(gaps_wb.sheetnames) == {"Technique Tracker"}

    assumptions_wb = load_workbook(io.BytesIO(build_xlsx_export(assessment, [], scope="assumptions")))
    assert set(assumptions_wb.sheetnames) == {"Not Applicable", "Assumptions"}


@pytest.mark.asyncio
async def test_xlsx_phase14h_polish(db_session):
    """Phase 14h: data bars + native bar chart on Coverage by Tactic, a
    3-color scale on the Priority column (replacing the old static per-cell
    fill), Read Me sheet protection, and workbook core properties that
    finally consume the branding override."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    wb = load_workbook(io.BytesIO(build_xlsx_export(
        assessment, [], branding={"report_display_name": "Acme Corp"},
    )))

    ws_tactic = wb["Coverage by Tactic"]
    assert len(ws_tactic.conditional_formatting._cf_rules) == 2  # Coverage % + Weighted %
    assert len(ws_tactic._charts) == 1

    ws_tracker = wb["Technique Tracker"]
    assert len(ws_tracker.conditional_formatting._cf_rules) == 1
    priority_cells = [c[7] for c in ws_tracker.iter_rows(min_row=2) if c[0].value]
    assert any(isinstance(c.value, int) for c in priority_cells)  # numeric, not "P2" string
    assert all(c.number_format == '"P"0' for c in priority_cells if isinstance(c.value, int))

    assert wb["Read Me"].protection.sheet is True

    assert assessment.name in wb.properties.title
    assert wb.properties.creator == "ScopeWise"
    assert "Acme Corp" in wb.properties.description


@pytest.mark.asyncio
async def test_html_report_pdf_metadata_tags(db_session):
    """Phase 14h: base.html carries <meta> tags WeasyPrint maps straight to
    PDF document properties (author -> /Author, description -> /Subject,
    keywords -> /Keywords; <title> -> /Title was already covered by the
    existing page_title assertions). Checked on the plain HTML string so
    this test runs without WeasyPrint's native libs."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    html = build_html_report(assessment, [], branding={"report_display_name": "Acme Corp"})
    assert '<meta name="author" content="Acme Corp">' in html
    assert '<meta name="description" content="MITRE ATT&amp;CK coverage assessment for' in html
    assert 'overall coverage 33.3%' in html
    assert '<meta name="keywords" content="MITRE ATT&amp;CK, coverage assessment, gap analysis, Acme Corp">' in html


def test_compare_golden():
    class A:  # minimal stand-in with the attributes compare reads
        pass

    baseline, current = A(), A()
    for a, strict, states in (
        (baseline, 40.0, {"T1059.001": "covered", "T1112": "not_covered",
                          "T1003.001": "covered", "T1200": "not_applicable",
                          "T1547.001": "covered"}),
        (current, 60.0, {"T1059.001": "covered", "T1112": "covered",
                         "T1003.001": "partial", "T1200": "covered",
                         "T1547.001": "not_applicable"}),
    ):
        a.assessment_id = uuid.uuid4()
        a.name = "x"
        a.completed_at = datetime.now(timezone.utc)
        a.attack_version = "19.1"
        a.summary = _summary(strict)
        a.technique_results = _results(states)

    diff = compare_assessments(current, baseline)
    assert [e["technique_id"] for e in diff["newly_covered"]] == ["T1112", "T1200"]
    assert [e["technique_id"] for e in diff["regressed"]] == ["T1003.001"]
    assert [e["technique_id"] for e in diff["na_changed"]] == ["T1200", "T1547.001"]
    assert diff["overall_delta"]["strict_pct"] == 20.0
    assert diff["attack_version_mismatch"] is False
    assert diff["tactic_deltas"] == [{
        "domain": "enterprise", "id": "TA0002", "name": "Execution",
        "current_strict_pct": 60.0, "baseline_strict_pct": 40.0, "delta": 20.0,
    }]
    # real technique names resolved from the pinned dataset
    assert diff["newly_covered"][0]["name"] == "Modify Registry"


def test_compare_degrades_on_malformed_jsonb():
    """technique_results/summary are unenforced JSONB — a row missing 'state'
    or a tactic missing 'id' (schema drift) must degrade, not 500
    (2026-08-02 adversarial review, non-blocking finding #1)."""
    class A:
        pass

    baseline, current = A(), A()
    for a in (baseline, current):
        a.assessment_id = uuid.uuid4()
        a.name = "x"
        a.completed_at = datetime.now(timezone.utc)
        a.attack_version = "19.1"
        a.summary = {"overall": {}, "domains": {"enterprise": {
            "tactics": [{"name": "no id here", "strict_pct": 10.0}]}}}
        a.technique_results = [
            {"technique_id": "T1059.001"},          # no 'state'
            {"technique_id": "T1112", "state": "covered"},
        ]
    # must not raise
    diff = compare_assessments(current, baseline)
    assert diff["tactic_deltas"] == []  # the id-less tactic is skipped, not crashed


# ---------------------------------------------------------------------------
# endpoint tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_report_endpoint_html_and_409(client, db_session):
    org, user, headers = await _make_user(db_session, role="viewer")  # viewers may read
    done = await _seed(db_session, org, user, uc_names=["Rule A"])
    pending = await _seed(db_session, org, user, status="pending")

    res = await client.get(
        f"/api/v1/mitre/assessments/{done.assessment_id}/report", headers=headers
    )
    assert res.status_code == 200
    assert res.json()["format"] == "html"
    assert "MITRE ATT&amp;CK Coverage Assessment" in res.json()["data"]

    res = await client.get(
        f"/api/v1/mitre/assessments/{pending.assessment_id}/report", headers=headers
    )
    assert res.status_code == 409


@pytest.mark.asyncio
async def test_trend_scopes_to_same_customer(client, db_session):
    """Phase A12: the auto-picked trend baseline must be the most recent
    completed run for the SAME customer, not just the same org."""
    org, user, headers = await _make_user(db_session, role="viewer")
    t0 = datetime.now(timezone.utc) - timedelta(days=2)

    acme_1 = await _seed(
        db_session, org, user, name="Acme Run 1", customer="Acme Corp",
        completed_at=t0,
    )
    acme_2 = await _seed(
        db_session, org, user, name="Acme Run 2", customer="Acme Corp",
        completed_at=t0 + timedelta(days=1),
    )

    res = await client.get(
        f"/api/v1/mitre/assessments/{acme_2.assessment_id}/report", headers=headers
    )
    assert res.status_code == 200
    html = res.json()["data"]
    assert "Trend vs your previous run" in html
    assert "Acme Run 1" in html
    assert acme_1.assessment_id  # baseline picked is acme_1, not some other org run


@pytest.mark.asyncio
async def test_trend_skips_cross_customer_run(client, db_session):
    """A different customer's earlier run must NOT be picked as the
    baseline — that would produce a nonsensical diff (real prod symptom)."""
    org, user, headers = await _make_user(db_session, role="viewer")
    t0 = datetime.now(timezone.utc) - timedelta(days=2)

    await _seed(
        db_session, org, user, name="Acme Run 1", customer="Acme Corp",
        completed_at=t0,
    )
    globex_1 = await _seed(
        db_session, org, user, name="Globex Run 1", customer="Globex Inc",
        completed_at=t0 + timedelta(days=1),
    )

    res = await client.get(
        f"/api/v1/mitre/assessments/{globex_1.assessment_id}/report", headers=headers
    )
    assert res.status_code == 200
    html = res.json()["data"]
    assert "Trend vs your previous run" not in html  # no prior Globex run exists


@pytest.mark.asyncio
async def test_trend_null_customer_still_matches_null(client, db_session):
    """Orgs that never set a customer keep pre-A12 behavior: NULL is
    NOT DISTINCT FROM NULL, so the plain most-recent-run pick still works."""
    org, user, headers = await _make_user(db_session, role="viewer")
    t0 = datetime.now(timezone.utc) - timedelta(days=2)

    await _seed(db_session, org, user, name="No-customer Run 1", completed_at=t0)
    run_2 = await _seed(
        db_session, org, user, name="No-customer Run 2",
        completed_at=t0 + timedelta(days=1),
    )

    res = await client.get(
        f"/api/v1/mitre/assessments/{run_2.assessment_id}/report", headers=headers
    )
    assert res.status_code == 200
    html = res.json()["data"]
    assert "Trend vs your previous run" in html
    assert "No-customer Run 1" in html


@pytest.mark.asyncio
async def test_pdf_endpoint(client, db_session, require_weasyprint):
    import base64

    org, user, headers = await _make_user(db_session)
    done = await _seed(db_session, org, user, uc_names=["Rule A"])
    res = await client.get(
        f"/api/v1/mitre/assessments/{done.assessment_id}/report",
        headers=headers, params={"format": "pdf"},
    )
    assert res.status_code == 200
    assert base64.b64decode(res.json()["data"]).startswith(b"%PDF")


@pytest.mark.asyncio
async def test_xlsx_endpoint_streams_with_content_type(client, db_session):
    org, user, headers = await _make_user(db_session, role="viewer")
    done = await _seed(db_session, org, user, uc_names=[FORMULA])
    res = await client.get(
        f"/api/v1/mitre/assessments/{done.assessment_id}/export.xlsx", headers=headers
    )
    assert res.status_code == 200
    assert res.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert "attachment" in res.headers["content-disposition"]
    ws = load_workbook(io.BytesIO(res.content))["Use-Case Mappings"]
    assert ws["B2"].value == "'" + FORMULA


@pytest.mark.asyncio
async def test_executive_scope_report(client, db_session):
    """scope=executive returns the 1-3 page leadership cut: cover + executive
    section only — no detailed section, appendices, TOC, or cross-refs."""
    org, user, headers = await _make_user(db_session, role="viewer")
    done = await _seed(db_session, org, user)
    res = await client.get(
        f"/api/v1/mitre/assessments/{done.assessment_id}/report",
        params={"format": "html", "scope": "executive"},
        headers=headers,
    )
    assert res.status_code == 200
    html = res.json()["data"]
    assert "Top 5 things to fix first" in html
    assert "Executive summary" in html
    assert "Appendix: technique register" not in html
    assert "Gap register" not in html
    assert "<h3>Contents</h3>" not in html
    assert "class='xref'" not in html
    # Phase A11 piece 3: no forced page breaks left in the trimmed
    # executive cut -- cover + executive summary flow continuously.
    assert 'class="page-break"' not in html
    res = await client.get(
        f"/api/v1/mitre/assessments/{done.assessment_id}/report",
        params={"scope": "bogus"},
        headers=headers,
    )
    assert res.status_code == 422


@pytest.mark.asyncio
async def test_rename_archive_and_list_filter(client, db_session):
    """Phase 14f: PATCH rename + soft-archive flag; archived rows leave the
    default list but return with include_archived=true. No delete exists."""
    org, user, headers = await _make_user(db_session, role="admin")
    a = await _seed(db_session, org, user)
    aid = str(a.assessment_id)

    res = await client.patch(
        f"/api/v1/mitre/assessments/{aid}", json={"name": "Renamed run"}, headers=headers
    )
    assert res.status_code == 200 and res.json()["name"] == "Renamed run"

    res = await client.patch(
        f"/api/v1/mitre/assessments/{aid}", json={"archived": True}, headers=headers
    )
    assert res.status_code == 200 and res.json()["archived"] is True

    res = await client.get("/api/v1/mitre/assessments", headers=headers)
    assert all(i["assessment_id"] != aid for i in res.json())

    res = await client.get(
        "/api/v1/mitre/assessments", params={"include_archived": True}, headers=headers
    )
    row = next(i for i in res.json() if i["assessment_id"] == aid)
    assert row["archived"] is True and row["name"] == "Renamed run"

    res = await client.patch(
        f"/api/v1/mitre/assessments/{aid}", json={}, headers=headers
    )
    assert res.status_code == 422


@pytest.mark.asyncio
async def test_compare_endpoint_and_authz(client, db_session):
    org, user, headers = await _make_user(db_session)
    baseline = await _seed(
        db_session, org, user,
        strict_pct=40.0,
        results=_results({"T1112": "not_covered", "T1059.001": "covered"}),
    )
    current = await _seed(
        db_session, org, user,
        strict_pct=60.0,
        results=_results({"T1112": "covered", "T1059.001": "covered"}),
    )
    pending = await _seed(db_session, org, user, status="pending")

    res = await client.get(
        f"/api/v1/mitre/assessments/{current.assessment_id}/compare/{baseline.assessment_id}",
        headers=headers,
    )
    assert res.status_code == 200
    body = res.json()
    assert [e["technique_id"] for e in body["newly_covered"]] == ["T1112"]
    assert body["regressed"] == [] and body["na_changed"] == []
    assert body["overall_delta"]["strict_pct"] == 20.0

    # baseline not completed -> 409
    res = await client.get(
        f"/api/v1/mitre/assessments/{current.assessment_id}/compare/{pending.assessment_id}",
        headers=headers,
    )
    assert res.status_code == 409

    # cross-org -> 404 on either side
    _, _, headers_b = await _make_user(db_session)
    res = await client.get(
        f"/api/v1/mitre/assessments/{current.assessment_id}/compare/{baseline.assessment_id}",
        headers=headers_b,
    )
    assert res.status_code == 404


@pytest.mark.asyncio
async def test_list_includes_domains_brief(client, db_session):
    org, user, headers = await _make_user(db_session)
    await _seed(db_session, org, user)
    res = await client.get("/api/v1/mitre/assessments", headers=headers)
    assert res.status_code == 200
    brief = res.json()[0]["domains_brief"]
    assert brief["enterprise"]["strict_pct"] == 33.3
    assert brief["enterprise"]["applicable"] == 3


@pytest.mark.asyncio
async def test_xlsx_summary_never_triggered_caveat(db_session):
    """2026-08-13: rules with Last Triggered 'never' add a source-health
    caveat under the roadmap bullet; without any, no caveat."""
    org, user, _ = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)
    uc = {"row_ref": "s:1", "name": "r1", "description": None, "logic": None,
          "log_source": None, "enabled": True, "mappings": [],
          "mapping_status": "customer_tagged"}
    wb = load_workbook(io.BytesIO(build_xlsx_export(
        assessment, [dict(uc, last_triggered="never"), dict(uc, row_ref="s:2")])))
    texts = [str(c.value) for row in wb["Summary"].iter_rows() for c in row if c.value]
    assert any("Caveat: 'buildable now'" in t and "1 of 2 rules" in t for t in texts)

    wb2 = load_workbook(io.BytesIO(build_xlsx_export(assessment, [uc])))
    texts2 = [str(c.value) for row in wb2["Summary"].iter_rows() for c in row if c.value]
    assert not any("Caveat: 'buildable now'" in t for t in texts2)


@pytest.mark.asyncio
async def test_pptx_builder_and_endpoint(client, db_session):
    """2026-08-14: PPTX briefing deck — builder returns a valid deck from the
    stored summary; endpoint streams it with the PowerPoint MIME type."""
    org, user, headers = await _make_user(db_session)
    assessment = await _seed(db_session, org, user)

    from pptx import Presentation as PptxPresentation

    from app.mitre.report import build_pptx_export

    data = build_pptx_export(assessment, [])
    assert data[:2] == b"PK"
    deck = PptxPresentation(io.BytesIO(data))
    assert len(deck.slides) >= 6
    texts = " ".join(
        sh.text_frame.text
        for s in deck.slides for sh in s.shapes if sh.has_text_frame
    )
    assert "MITRE ATT&CK" in texts
    assert "33.3%" in texts  # headline strict_pct from _summary()

    res = await client.get(
        f"/api/v1/mitre/assessments/{assessment.assessment_id}/export.pptx",
        headers=headers,
    )
    assert res.status_code == 200
    assert res.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml"
    )
    assert res.content[:2] == b"PK"
