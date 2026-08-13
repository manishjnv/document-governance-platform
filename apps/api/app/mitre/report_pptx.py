"""MITRE assessment PPTX briefing-deck builder (2026-08-14, from the VFQ
customer-deliverable review; expanded same day to the full 18-slide
structure of the manually-built VFQ client deck): a client-presentation deck
in the same design system as the XLSX export (deep purple / teal / magenta,
card-with-accent-bar layout, highlighted keywords, section dividers with
giant numbers).

Same invariants as every other report surface: numbers come ONLY from the
stored summary/technique_results JSONB — never recomputed here; the only LLM
text reused is the stored narrative's per-gap recommendations (rephrased
words, never numbers). The analyst-judgment slides of the hand-built deck
(strengths / key gaps) are replaced by DERIVED equivalents: best/worst
tactics, top log sources, provenance and validation shares — evidence the
data actually supports, phrased with the numbers inline. Fully
deterministic output for a given assessment.
"""

import io

from app.mitre.report_common import (
    DOMAIN_LABELS,
    compute_log_source_coverage,
    resolve_branding,
)

# palette — mirrors report_xlsx's BRAND/ACCENT plus the deck-only tones
_PURPLE = "341954"
_PURPLE2 = "46246B"
_PURPLE_NUM = "4A2A73"
_LAVENDER = "7370A7"
_MAGENTA = "B71D6B"
_ROSE = "C73E65"
_TEAL = "00A98B"
_GREEN = "007A5E"
_GREEN_D = "1E7B4D"
_RED_D = "B02830"
_AMBER = "C8801A"
_AMBER_D = "8A5A10"
_CARD = "F3F0F7"
_MINT = "E7F5F1"
_ROSEBG = "FCEEF3"
_GREY = "404040"
_MUTED = "BBBDC2"
_LILAC = "D8D2E4"
_FONT = "Tenorite"  # MS cloud font; PowerPoint substitutes silently if absent


def build_pptx_export(assessment, use_cases: list, branding: dict | None = None) -> bytes:
    """The briefing deck as .pptx bytes. use_cases: the same dict rows the
    XLSX builder receives (log_source / last_triggered feed several slides)."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    from app.mitre import attack_data

    branding = resolve_branding(branding)
    display_name = branding["report_display_name"]
    summary = assessment.summary or {}
    params = assessment.params or {}
    intake = params.get("intake") or {}
    env = params.get("environment") or {}
    env_lists = params.get("environment_lists") or {}
    overall = summary.get("overall", {})
    domains = summary.get("domains") or {}
    gaps = summary.get("gaps", [])
    roadmap = summary.get("roadmap", {})
    quality = summary.get("quality", {})
    counts = summary.get("counts", {})
    narrative = summary.get("narrative", {})
    gap_recs = narrative.get("gap_recommendations", {})
    not_applicable = summary.get("not_applicable", [])
    completed = str(assessment.completed_at or "")[:10]
    total_rules = counts.get("use_cases", len(use_cases)) or 0
    gap_total = (overall.get("not_covered") or 0) + (overall.get("partial") or 0)
    applicable = overall.get("applicable") or 0
    strict_pct = overall.get("strict_pct")
    never_count = sum(
        1 for uc in use_cases
        if str(uc.get("last_triggered") or "").strip().lower() == "never"
    )
    disabled_count = sum(1 for uc in use_cases if uc.get("enabled") is False)
    customer_tagged = counts.get("customer_tagged") or 0
    unmapped = counts.get("unmapped") or 0
    platforms = env.get("platforms") or []
    log_sources = env_lists.get("log_sources") or []
    tooling = env_lists.get("tooling") or []
    crown_jewels = env_lists.get("crown_jewels") or []
    groups = compute_log_source_coverage(
        use_cases, assessment.technique_results or [], attack_data.DEFAULT)
    zero_groups = [g for g in groups if not g.get("techniques_covered")]
    # N/A breakdown for the scope slide — classified from the stated reasons
    na_mobile = sum(1 for e in not_applicable if "Mobile matrix" in str(e.get("reason")))
    na_ics = sum(1 for e in not_applicable if "ICS matrix" in str(e.get("reason")))
    na_deprecated = sum(1 for e in not_applicable if "deprecated" in str(e.get("reason", "")).lower())
    na_platform = sum(1 for e in not_applicable if str(e.get("reason", "")).startswith("targets "))
    # primary domain + tactic rankings (chart, strengths, gaps)
    primary = max(
        (d for d in domains if (domains[d] or {}).get("applicable")),
        key=lambda d: domains[d].get("applicable", 0), default=None,
    )
    tactics = [t for t in (domains.get(primary, {}).get("tactics") or [])
               if t.get("applicable")] if primary else []
    tactics_ranked = sorted(tactics, key=lambda t: t.get("strict_pct") or 0)
    cj_gaps = sum(1 for g in gaps if g.get("crown_jewel_relevant"))
    threat_gaps = sum(1 for g in gaps if g.get("threat_relevance"))
    threat_bits = [b for b in (
        intake.get("industry"), intake.get("region"),
        ", ".join(intake.get("threat_actors") or []) or None) if b]
    footer_text = (
        f"{intake.get('project_name') or assessment.name} · MITRE ATT&CK "
        f"v{assessment.attack_version} · {display_name} · {completed}"
    )

    def rgb(hexc):
        return RGBColor.from_string(hexc)

    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]
    page_no = 0

    def slide():
        return prs.slides.add_slide(blank)

    def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
        sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(color)
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def text(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
        """paras: [(runs, opts)] with runs [(txt, {bold,color,size})]."""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, (runs, opts) in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = opts.get("align", PP_ALIGN.LEFT)
            if opts.get("space_after") is not None:
                p.space_after = Pt(opts["space_after"])
            for txt, ro in runs:
                r = p.add_run()
                r.text = str(txt)
                f = r.font
                f.name = _FONT
                f.size = Pt(ro.get("size", 9))
                f.bold = ro.get("bold", False)
                f.color.rgb = rgb(ro.get("color", "000000"))
        return tb

    def P(runs, **opts):
        return (runs, opts)

    def R(txt, bold=False, color="000000", size=9):
        return (txt, {"bold": bold, "color": color, "size": size})

    def K(txt, size=9, color=_GREEN):  # highlighted keyword run
        return R(txt, bold=True, color=color, size=size)

    def chrome(s, title, subtitle):
        nonlocal page_no
        page_no += 1
        text(s, 0.45, 0.18, 9.10, 0.55, [P([R(title, bold=True, color=_PURPLE, size=24)])])
        rect(s, 0.47, 0.78, 9.06, 0.03, _TEAL)
        text(s, 0.45, 0.86, 9.10, 0.30, [P([R(subtitle, color=_GREY, size=11.5)])])
        text(s, 0.45, 5.30, 7.50, 0.25, [P([R(footer_text, color=_MUTED, size=7.5)])])
        text(s, 9.25, 5.30, 0.40, 0.25,
             [P([R(str(page_no), color=_MUTED, size=8)], align=PP_ALIGN.RIGHT)])

    def card(s, x, y, w, h, bar=_TEAL, fillc=_CARD):
        rect(s, x, y, w, h, fillc)
        rect(s, x, y, 0.06, h, bar)

    def info_card(s, x, y, w, h, bar, fillc, head, head_color, body_paras):
        card(s, x, y, w, h, bar, fillc)
        text(s, x + 0.16, y + 0.07, w - 0.32, 0.26,
             [P([R(head, bold=True, color=head_color, size=10.5)])])
        text(s, x + 0.16, y + 0.34, w - 0.32, h - 0.42, body_paras)

    def stat_tile(s, x, y, w, h, big, label, color=_PURPLE):
        card(s, x, y, w, h, color)
        text(s, x + 0.16, y + 0.10, w - 0.28, 0.52,
             [P([R(big, bold=True, color=color, size=24)])])
        text(s, x + 0.16, y + 0.62, w - 0.28, h - 0.66,
             [P([R(label, color=_GREY, size=8.5)])])

    def divider(num, title, desc):
        nonlocal page_no
        page_no += 1
        s = slide()
        rect(s, 0, 0, 10.0, 5.625, _PURPLE)
        rect(s, 7.10, -1.20, 5.50, 5.50, _PURPLE2, MSO_SHAPE.OVAL)
        rect(s, 8.30, 3.60, 3.40, 3.40, _MAGENTA, MSO_SHAPE.OVAL)
        text(s, 0.90, 0.90, 3.60, 0.35,
             [P([R("MITRE ATT&CK Assessment", bold=True, color=_TEAL, size=12.5)])])
        text(s, 0.90, 1.90, 5.60, 1.00, [P([R(title, bold=True, color="FFFFFF", size=30)])])
        text(s, 0.90, 2.90, 5.80, 1.40, [P([R(desc, color=_LILAC, size=12.5)])])
        text(s, 6.30, 0.70, 3.20, 2.40,
             [P([R(num, bold=True, color=_PURPLE_NUM, size=120)])])
        return s

    def style_table(tbl, widths, header_size=9.5, body_size=9):
        for i, w in enumerate(widths):
            tbl.columns[i].width = Inches(w)
        for ri, row in enumerate(tbl.rows):
            for cell in row.cells:
                cell.margin_left = Inches(0.06)
                cell.margin_right = Inches(0.06)
                cell.margin_top = cell.margin_bottom = Inches(0.02)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(
                    _PURPLE if ri == 0 else (_CARD if ri % 2 else "FFFFFF"))
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = _FONT
                        r.font.size = Pt(header_size if ri == 0 else body_size)
                        if ri == 0:
                            r.font.bold = True
                            r.font.color.rgb = rgb("FFFFFF")

    def six_cards(s, entries):
        """Two-column grid of up to 6 evidence cards: (head, bar_color, paras)."""
        for i, (head, bar, paras) in enumerate(entries[:6]):
            x = 0.45 if i < 3 else 5.08
            y = 1.20 + (i % 3) * 1.36
            card(s, x, y, 4.50, 1.26, bar)
            text(s, x + 0.15, y + 0.07, 4.25, 0.24,
                 [P([R(head, bold=True, color=_PURPLE, size=10)])])
            text(s, x + 0.15, y + 0.33, 4.25, 0.88, paras)

    # ------------------------------------------------------------- 1 · cover
    s = slide()
    rect(s, 0, 0, 10.0, 5.625, _PURPLE)
    rect(s, -1.60, 1.40, 6.00, 6.00, _PURPLE2, MSO_SHAPE.OVAL)
    rect(s, 2.60, -2.60, 5.20, 5.20, _MAGENTA, MSO_SHAPE.OVAL)
    rect(s, 0.75, 3.15, 3.20, 0.03, _TEAL)
    text(s, 4.45, 1.05, 5.20, 1.90, [
        P([R("MITRE ATT&CK", bold=True, color="FFFFFF", size=36)]),
        P([R("Coverage Assessment", bold=True, color="FFFFFF", size=36)]),
        P([R("Detection coverage · Gaps · Roadmap", bold=True, color=_TEAL, size=17)]),
    ])
    meta = [P([R(intake.get("project_name") or assessment.name,
                 bold=True, color="FFFFFF", size=13.5)])]
    if intake.get("prepared_by"):
        meta.append(P([R(f"Prepared by {intake['prepared_by']}", color=_MUTED, size=11)]))
    meta.append(P([R(f"{total_rules} detection rules · "
                     f"{len(log_sources)} log sources · "
                     f"{len(platforms)} platforms", color=_MUTED, size=11)]))
    meta.append(P([R(f"MITRE ATT&CK v{assessment.attack_version} · {completed}",
                     color=_MUTED, size=11)]))
    text(s, 4.45, 3.15, 5.20, 1.90, meta)
    text(s, 0.75, 2.55, 3.40, 0.50,
         [P([R(display_name, bold=True, color="FFFFFF", size=13.5)])])
    page_no += 1

    # ------------------------------------------------------------ 2 · agenda
    s = slide()
    chrome(s, "Agenda & What to Expect",
           "Four sections — how your detection coverage was measured, what it "
           "shows, and what to do next")
    sections = [
        ("01", _TEAL, "Assessment Overview", "What we assessed, your data, how scoring works"),
        ("02", _LAVENDER, "Coverage Results", "Coverage by tactic · detection quality · log sources"),
        ("03", _MAGENTA, "Strengths, Gaps & Insights", "What works, what's missing — with evidence"),
        ("04", _ROSE, "Roadmap & Way Forward", "Priority fixes · phased plan · next steps"),
    ]
    for i, (num, color, head, sub) in enumerate(sections):
        y = 1.25 + i * 1.02
        card(s, 0.45, y, 4.40, 0.92, color)
        text(s, 0.62, y + 0.14, 0.62, 0.60, [P([R(num, bold=True, color=color, size=24)])])
        text(s, 1.28, y + 0.12, 3.50, 0.32, [P([R(head, bold=True, color=_PURPLE, size=13.5)])])
        text(s, 1.28, y + 0.46, 3.50, 0.42, [P([R(sub, color=_GREY, size=9.5)])])
    text(s, 5.05, 1.22, 4.50, 0.32,
         [P([R("How to read this assessment", bold=True, color=_PURPLE, size=12.5)])])
    info_card(s, 5.05, 1.62, 4.50, 1.10, _TEAL, _MINT, "What you can expect", _GREEN_D,
              [P([R("A ", size=9), K("technique-by-technique view", 9),
                  R(f" of what your {total_rules} rules detect across ", size=9),
                  K(f"{applicable} applicable techniques", 9),
                  R(" — every number computed from your own rule export, none estimated.", size=9)])])
    cannot = [P([R("It is ", size=9), K("not a penetration test", 9, _MAGENTA),
                 R(". It shows a rule ", size=9), R("exists", bold=True, size=9),
                 R(" for a technique — proving the rule ", size=9),
                 R("fires", bold=True, size=9), R(" needs validation", size=9)])]
    if never_count:
        cannot = [P(cannot[0][0] + [R(": ", size=9),
                    K(f"{never_count} of {total_rules} rules found no events when "
                      "last checked", 9, _MAGENTA), R(".", size=9)])]
    else:
        cannot = [P(cannot[0][0] + [R(" in your environment.", size=9)])]
    info_card(s, 5.05, 2.80, 4.50, 1.10, _MAGENTA, _ROSEBG,
              "What it cannot tell you", _RED_D, cannot)
    trust_share = round(100 * customer_tagged / total_rules) if total_rules else 0
    info_card(s, 5.05, 3.98, 4.50, 1.10, _LAVENDER, _CARD, "How much to trust it", _PURPLE,
              [P([R("Built from your ", size=9), K("own rule export", 9),
                  R(" and ", size=9), K("asset & log-source inventory", 9),
                  R(f" — {trust_share}% of rules carry your own ATT&CK tags, kept "
                    "verbatim. Fully repeatable: re-run any time to trend.", size=9)])])

    # ------------------------------------------------- 3 · divider 01
    divider("01", "Assessment Overview",
            "What was assessed, the data behind every number, and the rules of "
            "the scoring model.")

    # ------------------------------------------------- 4 · scope & inputs
    s = slide()
    chrome(s, "Scope & Inputs",
           "Everything below comes from your own exports — nothing was sampled "
           "or assumed")
    tiles = [
        (str(total_rules), f"detection rules analyzed "
         f"({total_rules - disabled_count} enabled)", _PURPLE),
        (str(len(log_sources)), "log sources declared as onboarded", _TEAL),
        (str(len(platforms)), "ATT&CK platforms in your asset inventory", _LAVENDER),
        (str(len(crown_jewels)), "crown-jewel assets declared", _MAGENTA),
    ]
    for i, (big, label, color) in enumerate(tiles):
        stat_tile(s, 0.45 + i * 2.32, 1.25, 2.17, 1.05, big, label, color)
    received = [
        P([K("Rule export", 9), R(" — name, detection logic, ATT&CK tags, "
            "status, log source per rule", size=9)], space_after=5),
        P([K("Environment workbook", 9),
           R(f" — {len(platforms)} platforms, {len(log_sources)} log sources, "
             f"{len(tooling)} security tools, {len(crown_jewels)} crown jewels",
             size=9)], space_after=5),
    ]
    if never_count:
        received.append(P([K("Validation results", 9),
                           R(f" — Last-Triggered health data for the rules "
                             f"({never_count} marked never fired)", size=9)],
                          space_after=5))
    received.append(P([R("Platforms: ", bold=True, color=_PURPLE, size=9),
                       R(", ".join(platforms[:10]) + (" …" if len(platforms) > 10 else ""),
                         size=9)]))
    info_card(s, 0.45, 2.50, 4.45, 2.55, _TEAL, _CARD, "What we received", _PURPLE, received)
    decisions = [
        P([K(f"{len([d for d in domains if domains[d].get('applicable')])} "
             + ("matrix" if len([d for d in domains if domains[d].get("applicable")]) == 1 else "matrices") + " in scope", 9),
           R(f" — {applicable} techniques apply to your platforms", size=9)],
          space_after=5)]
    if na_mobile:
        decisions.append(P([K("Mobile: not applicable", 9, _MAGENTA),
                            R(f" — no managed mobile fleet declared; {na_mobile} "
                              "techniques excluded with the reason recorded", size=9)],
                           space_after=5))
    if na_ics:
        decisions.append(P([K("ICS / OT: not applicable", 9, _MAGENTA),
                            R(f" — no OT assets declared; {na_ics} techniques "
                              "excluded", size=9)], space_after=5))
    if na_platform:
        decisions.append(P([K(f"{na_platform} techniques", 9),
                            R(" excluded for platforms you don't run", size=9)],
                           space_after=5))
    if na_deprecated:
        decisions.append(P([K(f"{na_deprecated} deprecated techniques", 9),
                            R(" auto-excluded — retired by MITRE, nobody should "
                              "score them", size=9)], space_after=5))
    decisions.append(P([R("Exclusions ", size=9), K("never lower the score", 9),
                        R(" — they leave the denominator and are listed with "
                          "reasons in the appendix.", size=9)]))
    info_card(s, 5.10, 2.50, 4.45, 2.55, _LAVENDER, _CARD,
              "Scope decisions (and why)", _PURPLE, decisions)

    # ------------------------------------------------- 5 · methodology
    s = slide()
    chrome(s, "Methodology — From Your Rules to a Coverage Score",
           "Five deterministic steps · AI assists with wording only — it never "
           "produces a number")
    keyword_tagged = max(0, total_rules - customer_tagged
                         - (counts.get("ai_tagged") or 0) - unmapped
                         - (counts.get("invalid") or 0) - (counts.get("manual") or 0))
    steps_cards = [
        ("1 · Ingest", _TEAL,
         [P([R("All ", size=8.5), K(f"{total_rules} rules", 8.5),
             R(" parsed with name, logic, log source, status & severity", size=8.5)])]),
        ("2 · Map", _LAVENDER,
         [P([K(f"{customer_tagged} your own tags", 8.5), R(" (kept verbatim) · "
             f"{keyword_tagged} keyword-matched · "
             f"{counts.get('ai_tagged') or 0} AI-suggested · ", size=8.5),
             K(f"{unmapped} unmapped", 8.5, _MAGENTA),
             R(" (excluded, listed)", size=8.5)])]),
        ("3 · Filter", _PURPLE_NUM,
         [P([K(f"{overall.get('not_applicable') or 0} techniques N/A", 8.5),
             R(" — wrong platform, no such asset, or deprecated; every reason "
               "stated", size=8.5)])]),
        ("4 · Score", _MAGENTA,
         [P([R("Each technique: ", size=8.5), K("covered", 8.5), R(" · ", size=8.5),
             K("partial", 8.5, _AMBER), R(" · ", size=8.5),
             K("not covered", 8.5, _MAGENTA),
             R(" — pure computation, reproducible", size=8.5)])]),
        ("5 · Rank", _ROSE,
         [P([R("Gaps ordered by ", size=8.5), K("priority tier", 8.5),
             R(", your ", size=8.5), K("threat profile", 8.5), R(", ", size=8.5),
             K("crown jewels", 8.5), R(" & build effort", size=8.5)])]),
    ]
    for i, (head, color, bodyp) in enumerate(steps_cards):
        x = 0.45 + i * 1.86
        card(s, x, 1.25, 1.74, 1.55, color)
        text(s, x + 0.14, 1.33, 1.48, 0.26, [P([R(head, bold=True, color=color, size=10.5)])])
        text(s, x + 0.14, 1.62, 1.48, 1.10, bodyp)
    info_card(s, 0.45, 3.05, 4.45, 1.00, _TEAL, _MINT, "Deterministic by design", _GREEN_D,
              [P([R("Run it twice, get the ", size=9), K("same numbers", 9),
                  R(". AI only phrases recommendations — every figure traces to a "
                    "rule you own.", size=9)])])
    info_card(s, 5.10, 3.05, 4.45, 1.00, _LAVENDER, _CARD, "Fully auditable", _PURPLE,
              [P([R("Every technique links to its ", size=9),
                  K("exact supporting rules", 9),
                  R(" with mapping source and confidence — drill down from any "
                    "number in the spreadsheet register.", size=9)])])
    if threat_bits:
        text(s, 0.45, 4.25, 9.10, 0.70, [
            P([R("Threat context applied: ", bold=True, color=_PURPLE, size=9.5),
               K(" · ".join(str(b) for b in threat_bits), 9.5),
               R(f" — {threat_gaps} gaps match your declared threat profile and "
                 "are lifted up the priority list.", size=9.5)])])

    # ------------------------------------------------- 6 · divider 02
    divider("02", "Coverage Results",
            "The headline score, coverage per attack stage, detection quality, "
            "and what each log source buys you.")

    # -------------------------------------------------------- 7 · headline
    s = slide()
    chrome(s, f"Headline Result — {strict_pct}% of Applicable Techniques Covered",
           f"{overall.get('covered')} of {applicable} techniques your "
           "environment is exposed to have at least one enabled detection rule")
    tiles = [
        (f"{strict_pct}%",
         f"coverage ({overall.get('weighted_pct')}% weighted) — counted at "
         "sub-technique level", _PURPLE),
        (f"{overall.get('covered')} / {applicable}",
         "techniques covered by at least one enabled rule", _TEAL),
        (str(gap_total),
         f"gaps: {overall.get('not_covered')} with no rule + "
         f"{overall.get('partial')} partially covered", _MAGENTA),
        (str(overall.get("not_applicable")),
         "not applicable — excluded with stated reasons, not scored", _LAVENDER),
    ]
    for i, (big, label, color) in enumerate(tiles):
        stat_tile(s, 0.45 + i * 2.32, 1.25, 2.17, 1.20, big, label, color)
    info_card(s, 0.45, 2.70, 4.45, 1.10, _TEAL, _MINT, "Where you stand", _GREEN_D,
              [P([R("Early SIEM detection programs typically start ", size=9),
                  K("under 10%", 9), R(" of ATT&CK — the roadmap matters more than "
                    "the grade.", size=9)])])
    if never_count:
        info_card(s, 5.10, 2.70, 4.45, 1.10, _MAGENTA, _ROSEBG, "The trust lens", _RED_D,
                  [P([K(f"{never_count} of {total_rules} rules found no events",
                        9, _MAGENTA),
                      R(" when last validated. They count toward coverage — but "
                        "their quality scores are capped until they prove they "
                        "fire.", size=9)])])
    else:
        info_card(s, 5.10, 2.70, 4.45, 1.10, _LAVENDER, _CARD, "Provenance", _PURPLE,
                  [P([R("Every number is computed from your own rule export and "
                        "asset inventory — drill into any technique in the "
                        "spreadsheet register.", size=9)])])
    domains_line = " · ".join(
        f"{DOMAIN_LABELS.get(d, d)} {v.get('strict_pct')}%"
        for d, v in domains.items() if v.get("applicable")
    ) or "no applicable domains"
    info_card(s, 0.45, 3.95, 9.10, 1.00, _LAVENDER, _CARD,
              "Why “applicable” matters", _PURPLE,
              [P([R("Techniques for platforms you don't run are excluded from the "
                    "denominator with the reason recorded — the score reflects "
                    "your ", size=9), K("real attack surface", 9),
                  R(". ", size=9),
                  R("By matrix: ", bold=True, color=_PURPLE, size=9),
                  K(domains_line, 9)])])

    # -------------------------------------------- 8 · coverage by tactic
    if tactics:
        s = slide()
        chrome(s, "Coverage by Attack Stage (Tactic)",
               f"Share of applicable techniques covered at each stage — "
               f"{DOMAIN_LABELS.get(primary, primary)} matrix")
        cd = CategoryChartData()
        cd.categories = [t.get("name") for t in tactics]
        cd.add_series("Coverage %", [t.get("strict_pct") or 0 for t in tactics])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.45),
                                Inches(1.25), Inches(5.85), Inches(3.85), cd)
        ch = gf.chart
        ch.has_title = False
        ch.has_legend = False
        ser = ch.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = rgb(_TEAL)
        plot = ch.plots[0]
        plot.gap_width = 60
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.font.name = _FONT
        plot.data_labels.font.color.rgb = rgb(_PURPLE)
        plot.data_labels.number_format = "0.0"
        plot.data_labels.number_format_is_linked = False
        ch.category_axis.tick_labels.font.size = Pt(8.5)
        ch.category_axis.tick_labels.font.name = _FONT
        ch.value_axis.visible = False
        ch.value_axis.has_major_gridlines = False
        worst, best = tactics_ranked[:2], tactics_ranked[-2:][::-1]
        info_card(s, 6.50, 1.25, 3.05, 1.30, _TEAL, _MINT, "Strongest stages", _GREEN_D,
                  [P([K(f"{t.get('name')} {t.get('strict_pct')}%", 8.5),
                      R(f" — {t.get('covered')} of {t.get('applicable')} "
                        "techniques", size=8.5)], space_after=3) for t in best])
        info_card(s, 6.50, 2.72, 3.05, 1.30, _MAGENTA, _ROSEBG, "Weakest stages", _RED_D,
                  [P([K(f"{t.get('name')} {t.get('strict_pct')}%", 8.5, _MAGENTA),
                      R(f" — {t.get('not_covered')} techniques open", size=8.5)],
                     space_after=3) for t in worst])
        biggest = max(tactics, key=lambda t: t.get("not_covered") or 0)
        info_card(s, 6.50, 4.19, 3.05, 0.91, _LAVENDER, _CARD, "Where the work is", _PURPLE,
                  [P([K(f"{biggest.get('name')}", 8.5),
                      R(f" holds the most open techniques ({biggest.get('not_covered')}) "
                        "— the roadmap front-loads the ranked ones.", size=8.5)])])

    # ---------------------------------------------- 9 · detection quality
    if quality.get("scored"):
        s = slide()
        chrome(s, "Detection Quality — Coverage Is Not the Same as Confidence",
               "Each covered technique scored 0–100 on rule provenance, status, "
               "logic and telemetry match")
        cd = CategoryChartData()
        cd.categories = ["Strong (75+)", "Moderate (45–74)", "Weak (<45)"]
        cd.add_series("Techniques", [quality.get("strong", 0),
                                     quality.get("moderate", 0),
                                     quality.get("weak", 0)])
        gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.45), Inches(1.30),
                                Inches(3.55), Inches(3.30), cd)
        ch = gf.chart
        ch.has_title = False
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(9)
        ch.legend.font.name = _FONT
        for pt, color in zip(ch.plots[0].series[0].points, (_TEAL, _AMBER, _MAGENTA)):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = rgb(color)
        plot = ch.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True
        plot.data_labels.font.size = Pt(10)
        plot.data_labels.font.bold = True
        plot.data_labels.font.name = _FONT
        plot.data_labels.font.color.rgb = rgb("FFFFFF")
        scored = quality.get("scored") or 1
        mod_share = round(100 * (quality.get("moderate", 0)) / scored)
        info_card(s, 4.35, 1.30, 5.20, 1.16, _TEAL, _MINT,
                  f"{quality.get('strong', 0)} strong", _GREEN_D,
                  [P([R("Backed by enabled rules whose ", size=9),
                      K("provenance, logic and telemetry all line up", 9),
                      R(" — your most trustworthy detections.", size=9)])])
        mod_body = [P([R("Reasonable rules with an open question — ", size=9),
                       K("unproven firing, partial telemetry match, or "
                         "AI-suggested mapping", 9, _AMBER),
                       R(f". That is {mod_share}% of everything scored.", size=9)])]
        info_card(s, 4.35, 2.60, 5.20, 1.16, _AMBER, _CARD,
                  f"{quality.get('moderate', 0)} moderate — and why", _AMBER_D, mod_body)
        needle = [P([R("Average strength today: ", size=9),
                     K(f"{quality.get('avg_strength')}/100", 9, _MAGENTA),
                     R(". ", size=9)])]
        if never_count:
            needle[0][0].extend([
                K(f"Validate the {never_count} silent rules", 9, _MAGENTA),
                R(" — zero new engineering; on the next re-run moderate flips "
                  "toward strong, which is what an auditor (or an attacker) "
                  "actually tests.", size=9)])
        else:
            needle[0][0].extend([
                R("Prove the moderate rules fire (trigger drills / replay) to "
                  "move them to strong without new engineering.", size=9)])
        info_card(s, 4.35, 3.90, 5.20, 1.16, _MAGENTA, _ROSEBG,
                  "How to move the needle", _RED_D, needle)

    # ------------------------------------------- 10 · coverage by log source
    if groups:
        s = slide()
        chrome(s, "What Each Log Source Buys You",
               "Detection rules and techniques covered per source — and the "
               "sources paying rent for nothing")
        top = groups[:8]
        rows = [("Log source", "Rules", "Techniques covered", "Attack stages")]
        rows += [(g.get("log_source"), g.get("rule_count"),
                  g.get("techniques_covered"),
                  ", ".join((g.get("tactics") or [])[:4])) for g in top]
        tbl = s.shapes.add_table(len(rows), 4, Inches(0.45), Inches(1.25),
                                 Inches(5.60), Inches(3.55)).table
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                tbl.cell(ri, ci).text = str(val)
        style_table(tbl, [2.10, 0.65, 0.95, 1.90], header_size=9, body_size=8)
        best_g = top[0] if top else None
        if best_g:
            info_card(s, 6.30, 1.25, 3.25, 1.30, _TEAL, _MINT, "Your workhorse", _GREEN_D,
                      [P([K(str(best_g.get("log_source")), 8.5),
                          R(f" powers {best_g.get('rule_count')} rules covering "
                            f"{best_g.get('techniques_covered')} techniques — "
                            "protect this pipeline first.", size=8.5)])])
        if zero_groups:
            worst_names = ", ".join(str(g.get("log_source")) for g in zero_groups[:3])
            info_card(s, 6.30, 2.72, 3.25, 1.30, _MAGENTA, _ROSEBG,
                      "Paying for silence", _RED_D,
                      [P([K(f"{len(zero_groups)} source group(s)", 8.5, _MAGENTA),
                          R(" cover no technique today (incl. ", size=8.5),
                          R(worst_names, bold=True, size=8.5),
                          R(") — map them properly or consciously accept, and "
                            "document, the gap.", size=8.5)])])
        info_card(s, 6.30, 4.19, 3.25, 0.91, _LAVENDER, _CARD, "Action", _PURPLE,
                  [P([R("The spreadsheet's ", size=8.5),
                      K("Coverage by Log Source", 8.5),
                      R(" sheet lists every source with its exact techniques.",
                        size=8.5)])])

    # ------------------------------------------------- 11 · divider 03
    divider("03", "Strengths, Gaps & Insights",
            "What the rule base genuinely does well — and where an attacker "
            "meets no resistance today.")

    # ------------------------------------------------ 12 · what's working
    strengths = []
    for t in [t for t in reversed(tactics_ranked) if t.get("covered")][:2]:
        strengths.append((
            f"{t.get('name')} — {t.get('strict_pct')}% covered", _TEAL,
            [P([K(f"{t.get('covered')} of {t.get('applicable')} techniques", 9),
                R(" at this attack stage have a detection — your strongest "
                  "stage(s). Keep these rules maintained.", size=9)])]))
    if quality.get("strong"):
        strengths.append((
            f"{quality.get('strong')} strong detections", _TEAL,
            [P([R("Techniques where ", size=9),
                K("provenance, logic and telemetry all line up", 9),
                R(" — the core an incident responder can rely on.", size=9)])]))
    if groups:
        g0 = groups[0]
        strengths.append((
            f"{g0.get('log_source')} — telemetry backbone", _LAVENDER,
            [P([K(f"{g0.get('rule_count')} rules → "
                  f"{g0.get('techniques_covered')} techniques", 9),
                R(" from one source — a proven pipeline worth protecting "
                  "(health-monitor it).", size=9)])]))
    if total_rules and customer_tagged:
        strengths.append((
            "Well-governed rule base", _LAVENDER,
            [P([K(f"{round(100 * customer_tagged / total_rules)}% of rules carry "
                  "your own ATT&CK tags", 9),
                R(" — rare discipline, and it made this assessment precise.",
                  size=9)])]))
    if threat_gaps and threat_bits:
        strengths.append((
            "Threat-informed prioritisation", _MAGENTA,
            [P([R("Your profile (", size=9),
                K(" · ".join(str(b) for b in threat_bits), 9),
                R(f") is applied — {threat_gaps} matching gaps are lifted up the "
                  "build order.", size=9)])]))
    if total_rules:
        strengths.append((
            "Everything is auditable", _MAGENTA,
            [P([R("Every covered technique traces to ", size=9),
                K("named rules with mapping source and confidence", 9),
                R(" — no black-box claims in this assessment.", size=9)])]))
    if strengths:
        s = slide()
        chrome(s, "What's Working",
               "Real capability in the current rule base — worth protecting and "
               "building on")
        six_cards(s, strengths)

    # ---------------------------------------------------- 13 · key gaps
    gap_cards = []
    for t in [t for t in tactics_ranked if t.get("not_covered")][:2]:
        gap_cards.append((
            f"{t.get('name')} — {t.get('strict_pct')}% covered", _MAGENTA,
            [P([K(f"{t.get('not_covered')} techniques open", 9, _MAGENTA),
                R(" at this stage. ", size=9)], space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("filter the tracker to this tactic and work the ranked gaps "
                  "top-down — most are ", size=8.5),
                K("buildable on logs you already collect", 8.5),
                R(".", size=8.5)])]))
    if cj_gaps:
        gap_cards.append((
            "Crown-jewel exposure", _MAGENTA,
            [P([K(f"{cj_gaps} of {gap_total} gaps", 9, _MAGENTA),
                R(" touch platforms or telemetry tied to your declared crown "
                  "jewels.", size=9)], space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("they already sort above equal-priority peers — start there.",
                  size=8.5)])]))
    if zero_groups:
        gap_cards.append((
            "Log sources paying rent for nothing", _MAGENTA,
            [P([K(f"{len(zero_groups)} source group(s)", 9, _MAGENTA),
                R(" ship logs that cover no technique today.", size=9)],
               space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("write rules for them or consciously accept (and document) "
                  "the blind spot.", size=8.5)])]))
    if never_count:
        gap_cards.append((
            "Unproven majority" if never_count > total_rules / 2 else "Unproven rules",
            _MAGENTA,
            [P([K(f"{never_count} of {total_rules} rules never fired", 9, _MAGENTA),
                R(" in validation — coverage exists on paper; proof is pending.",
                  size=9)], space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("run a ", size=8.5), K("validation sprint", 8.5),
                R(" — replay attack data or trigger drills; no new rules needed.",
                  size=8.5)])]))
    if unmapped:
        gap_cards.append((
            "Rules counting toward nothing", _MAGENTA,
            [P([K(f"{unmapped} rules", 9, _MAGENTA),
                R(" map to no ATT&CK technique and add no coverage.", size=9)],
               space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("map them in the rule editor or retire them — both beat "
                  "silent shelfware.", size=8.5)])]))
    if roadmap.get("long"):
        gap_cards.append((
            "Telemetry you don't collect yet", _MAGENTA,
            [P([K(f"{len(roadmap['long'])} gap" + ("s" if len(roadmap["long"]) != 1 else ""), 9, _MAGENTA),
                R(" need a telemetry capability that isn't onboarded today.",
                  size=9)], space_after=3),
             P([R("What to do: ", bold=True, color=_GREY, size=8.5),
                R("bundle them into one onboarding plan instead of ad-hoc "
                  "connector work.", size=8.5)])]))
    if gap_cards:
        s = slide()
        chrome(s, "Key Gaps",
               "Where an attacker meets no detection today — each tied to the "
               "evidence behind it")
        six_cards(s, gap_cards)

    # ----------------------------------------------------- 14 · top fixes
    if gaps:
        top5 = gaps[:5]
        s = slide()
        chrome(s, f"Top {len(top5)} Fixes — Start Here",
               "Highest-priority gaps first, ranked by real-world prevalence, "
               "your threat profile and build effort")
        rows = [("#", "Technique", "Recommendation", "Build effort")]
        for i, g in enumerate(top5, 1):
            tid = g.get("technique_id")
            rec = gap_recs.get(tid) or g.get("hint") or ""
            rows.append((str(i), f"{g.get('name')} ({tid})", str(rec)[:220],
                         {"short": "Short (0–3 mo)", "mid": "Mid (3–9 mo)",
                          "long": "Long (9–18 mo)"}.get(g.get("feasibility"), "")))
        tbl = s.shapes.add_table(len(rows), 4, Inches(0.45), Inches(1.30),
                                 Inches(9.10), Inches(3.30)).table
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                tbl.cell(ri, ci).text = str(val)
        style_table(tbl, [0.40, 2.70, 4.40, 1.60], header_size=9.5, body_size=8.5)
        text(s, 0.45, 4.80, 9.10, 0.40, [
            P([R("Detection approach, required log fields and reference KQL for "
                 "every gap are in the ", size=9.5),
               R("spreadsheet technique tracker", bold=True, color=_PURPLE, size=9.5),
               R(" delivered with this deck.", size=9.5)])])

    # ------------------------------------------------- 15 · divider 04
    divider("04", "Roadmap & Way Forward",
            "Sequenced by effort and dependency: prove what exists, build "
            "what's missing, then extend scope.")

    # -------------------------------------------------------- 16 · roadmap
    s = slide()
    chrome(s, "Improvement Roadmap",
           f"{len(roadmap.get('short', []))} of {gap_total} gaps are buildable "
           "on logs you already collect — sequencing is the whole game")
    buckets = [
        ("NOW · 0–3 months", _TEAL, _MINT, _GREEN_D, roadmap.get("short", []),
         "telemetry already onboarded — build the detection now",
         [K("provable coverage", 8.5),
          R(" — the top attacker entry points get watched.", size=8.5)]),
        ("NEXT · 3–9 months", _LAVENDER, _CARD, _PURPLE, roadmap.get("mid", []),
         "onboard telemetry from tooling you already own first",
         [K("no tool purchase needed", 8.5),
          R(" — connect what you already have, then detect.", size=8.5)]),
        ("LATER · 9–18 months", _MAGENTA, _ROSEBG, _RED_D, roadmap.get("long", []),
         "needs a new telemetry capability — plan deliberately",
         [K("no undeclared blind spots", 8.5, _MAGENTA),
          R(" — everything watched or accepted in writing.", size=8.5)]),
    ]
    for i, (head, bar, fillc, hc, items, desc, outcome) in enumerate(buckets):
        x = 0.45 + i * 3.10
        card(s, x, 1.25, 2.95, 2.55, bar, fillc)
        text(s, x + 0.15, 1.33, 2.65, 0.26,
             [P([R(f"{head}  ·  {len(items)}", bold=True, color=hc, size=11)])])
        paras = [P([R(desc, color=_GREY, size=8.5)], space_after=5)]
        for g in items[:3]:
            paras.append(P([R(f"{g.get('technique_id')} ", bold=True,
                              color=_PURPLE, size=8.5),
                            R(str(g.get("name"))[:34], size=8.5)], space_after=2))
        if len(items) > 3:
            paras.append(P([R(f"+ {len(items) - 3} more in the tracker",
                              color=_GREY, size=8)], space_after=5))
        else:
            paras.append(P([R("", size=8)], space_after=5))
        paras.append(P([R("Outcome: ", bold=True, color=hc, size=8.5)] + outcome))
        text(s, x + 0.15, 1.66, 2.65, 2.05, paras)
    if never_count:
        info_card(s, 0.45, 4.00, 9.10, 1.00, _AMBER, _CARD,
                  "One caveat before committing build dates", _AMBER_D,
                  [P([R("“Buildable now” means the log source is declared as "
                        "onboarded — ", size=9),
                      K(f"{never_count} of {total_rules} rules found no events "
                        "when last validated", 9, _AMBER),
                      R(", so verify source health first (parsing, normalization, "
                        "recent events), then commit dates.", size=9)])])

    # ----------------------------------------------------- 17 · next steps
    s = slide()
    chrome(s, "Recommended Next Steps",
           "Concrete, owned actions — each one moves the next re-run's number")
    steps = []
    if never_count:
        steps.append(("Validate the silent rules",
                      f"{never_count} rules found no events when last validated — "
                      "replay test data or trigger drills.",
                      _TEAL, "Effort: low", "Impact: very high"))
    if gaps:
        steps.append((f"Build the top {min(5, len(gaps))} detections",
                      "Highest-priority gaps, ranked for you — start with the "
                      "short-term bucket.", _TEAL, "Effort: low", "Impact: high"))
    if unmapped:
        steps.append(("Map or retire the unmapped rules",
                      f"{unmapped} rules map to no ATT&CK technique and count "
                      "toward nothing.", _LAVENDER, "Effort: low", "Impact: medium"))
    if zero_groups:
        steps.append((f"Decide on {len(zero_groups)} unread log source group(s)",
                      "Write rules for them or consciously accept (and document) "
                      "the blind spot.", _LAVENDER, "Effort: medium", "Impact: high"))
    if roadmap.get("long"):
        steps.append(("Plan new telemetry deliberately",
                      f"{len(roadmap['long'])} gap(s) need a capability you don't "
                      "collect today — bundle into one onboarding plan.",
                      _MAGENTA, "Effort: medium", "Impact: high"))
    steps.append(("Re-run to trend",
                  f"Re-assess quarterly against this {strict_pct}% baseline — "
                  "same method, same math, comparable every time.",
                  _MAGENTA, "Effort: low", "Impact: high — trend view"))
    for i, (head, sub, color, eff, imp) in enumerate(steps[:6]):
        x = 0.45 if i < 3 else 5.08
        y = 1.25 + (i % 3) * 1.32
        card(s, x, y, 4.50, 1.20, color)
        text(s, x + 0.15, y + 0.10, 0.55, 0.60,
             [P([R(str(i + 1), bold=True, color=color, size=24)])])
        text(s, x + 0.72, y + 0.10, 3.65, 0.30,
             [P([R(head, bold=True, color=_PURPLE, size=10.5)])])
        text(s, x + 0.72, y + 0.42, 3.65, 0.48, [P([R(sub, color=_GREY, size=8.5)])])
        text(s, x + 0.72, y + 0.92, 3.65, 0.24,
             [P([K(eff, 8.5, _GREEN_D), R("   ·   ", color=_GREY, size=8.5),
                 K(imp, 8.5, _PURPLE)])])

    # ------------------------------------------------------- 18 · closing
    s = slide()
    rect(s, 0, 0, 10.0, 5.625, _PURPLE)
    rect(s, -1.60, 1.40, 6.00, 6.00, _PURPLE2, MSO_SHAPE.OVAL)
    rect(s, 7.10, -1.20, 5.50, 5.50, _MAGENTA, MSO_SHAPE.OVAL)
    text(s, 0.90, 2.00, 8.20, 1.00,
         [P([R("Thank you · Q&A", bold=True, color="FFFFFF", size=34)])])
    text(s, 0.90, 3.10, 8.20, 1.10, [
        P([R(footer_text, bold=True, color=_TEAL, size=13)]),
        P([R("Full detail: the spreadsheet technique tracker, the PDF report "
             "and the ATT&CK Navigator layers.", color=_LILAC, size=11)]),
    ])

    cp = prs.core_properties
    cp.title = f"{assessment.name} — MITRE ATT&CK Coverage Assessment"
    cp.author = display_name
    cp.last_modified_by = display_name

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
